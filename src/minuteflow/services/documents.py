from __future__ import annotations

import csv
import json
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from minuteflow.models import DocumentParseResult
from minuteflow.utils import collect_texts, normalize_path


class DocumentService:
    def parse_document(self, path: str) -> dict:
        source = normalize_path(path)
        if not source.exists():
            raise FileNotFoundError(f"Document not found: {source}")
        suffix = source.suffix.lower()
        if suffix in {".md", ".markdown", ".txt"}:
            text = source.read_text(encoding="utf-8")
            backend = "plain-text"
        elif suffix == ".json":
            data = json.loads(source.read_text(encoding="utf-8"))
            text = json.dumps(data, ensure_ascii=False, indent=2)
            backend = "json"
        elif suffix == ".csv":
            text = self._parse_csv(source)
            backend = "csv"
        elif suffix == ".docx":
            text = self._parse_docx(source)
            backend = "python-docx"
        elif suffix == ".pptx":
            text = self._parse_pptx(source)
            backend = "python-pptx"
        elif suffix == ".pdf":
            text = self._parse_pdf(source)
            backend = "pypdf"
        elif suffix == ".xlsx":
            text = self._parse_xlsx(source)
            backend = "openpyxl"
        else:
            raise ValueError(f"Unsupported document type: {source.suffix}")

        result = DocumentParseResult(
            source_path=str(source),
            file_type=suffix.lstrip("."),
            title=source.stem,
            text=text.strip(),
            backend=backend,
        )
        return result.model_dump()

    def parse_documents(self, paths: list[str]) -> dict:
        results = [self.parse_document(path) for path in paths]
        return {
            "documents": results,
            "count": len(results),
        }

    def _parse_docx(self, path: Path) -> str:
        document = DocxDocument(path)
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        tables: list[str] = []
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    tables.append(" | ".join(cells))
        return collect_texts([*paragraphs, *tables])

    def _parse_pptx(self, path: Path) -> str:
        presentation = Presentation(path)
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            fragments = [f"# Slide {index}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    fragments.append(shape.text.strip())
            notes = self._extract_slide_notes(slide)
            if notes:
                fragments.append("Notes:")
                fragments.append(notes)
            slides.append("\n".join(fragments))
        return "\n\n".join(slides)

    def _extract_slide_notes(self, slide) -> str:
        notes_text: list[str] = []
        notes_slide = getattr(slide, "notes_slide", None)
        if not notes_slide:
            return ""
        for shape in notes_slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                notes_text.append(shape.text.strip())
        return "\n".join(notes_text)

    def _parse_pdf(self, path: Path) -> str:
        reader = PdfReader(path)
        pages = []
        for index, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                pages.append(f"# Page {index}\n{text}")
        return "\n\n".join(pages)

    def _parse_csv(self, path: Path) -> str:
        rows: list[str] = []
        with path.open("r", encoding="utf-8", newline="") as handle:
            reader = csv.reader(handle)
            for row in reader:
                if any(cell.strip() for cell in row):
                    rows.append(" | ".join(cell.strip() for cell in row))
        return "\n".join(rows)

    def _parse_xlsx(self, path: Path) -> str:
        workbook = load_workbook(path, data_only=True)
        sheets: list[str] = []
        for sheet in workbook.worksheets:
            lines = [f"# Sheet {sheet.title}"]
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value).strip() for value in row]
                if any(values):
                    lines.append(" | ".join(values))
            sheets.append("\n".join(lines))
        return "\n\n".join(sheets)

